# coding=utf-8

import logging
from mammut.common.storage.storage_base import SlidesStorageBase
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import namedtuple

logger = logging.getLogger()

SlideElement = namedtuple("SlideElement", ["slideNumber", "elementType", "content"])


class EmbeddedStoragePowerPoint(SlidesStorageBase):
    slides_metadata_re_string = (
        "^(sheet-id|Sheet-id|SHEET-ID)[ ]*->[ ]*(.+\\.(xlsx|pptx))$"
    )

    def __init__(self, storage_path):
        SlidesStorageBase.__init__(self, storage_path)

    def get_slide_elements(self, presentation_id):
        prs = Presentation(self.storage_path + presentation_id)
        slide_elements = []
        for slide in prs.slides:
            for j, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    url = shape.click_action.hyperlink.address
                    slide_elements.append((SlideElement(j + 1, "image", url)))
                elif shape.has_text_frame:
                    text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = text + run.text

                    slide_elements.append((SlideElement(j + 1, "text", text),))

        return slide_elements
